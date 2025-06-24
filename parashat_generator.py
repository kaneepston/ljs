import requests
import re
import html
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import datetime

def clean_text(raw_text):
    """
    A more robust function to remove all HTML tags, entities, and footnote content.
    """
    # Validate input
    if not isinstance(raw_text, str):
        print(f"Warning: clean_text received non-string input: {raw_text} (type: {type(raw_text)})")
        return str(raw_text) if raw_text is not None else ""
    
    # 1. Specifically target and remove the common footnote structure first.
    text = re.sub(r'<sup class="footnote-marker">.*?</sup><i class="footnote">.*?</i>', '', raw_text, flags=re.DOTALL)
    text = re.sub(r'<i class="footnote">.*?</i>', '', text, flags=re.DOTALL)
    
    # 2. Then, remove any other remaining HTML tags (like sup, b, span, br).
    text = re.sub(r'<.*?>', '', text)
    
    # 3. Un-escape any lingering HTML entities like &nbsp; or &thinsp;
    text = html.unescape(text)
    
    # 4. Remove any standalone footnote markers that might be left
    text = text.replace('*', '')

    # 5. Clean up any strange artifacts that might result from cleaning, like double commas.
    text = text.replace(',,', ',')

    # 6. Return the text with leading/trailing whitespace removed.
    return text.strip()

def get_next_shabbat_date(weeks_ahead=0):
    """
    Calculate the date of the next Shabbat (Saturday) that is 'weeks_ahead' weeks from now.
    """
    today = datetime.datetime.now()
    
    # Find the next Saturday (Shabbat)
    # Monday = 0, Tuesday = 1, ..., Saturday = 5, Sunday = 6
    days_until_saturday = (5 - today.weekday()) % 7
    
    # If today is Saturday, days_until_saturday will be 0, so we want next Saturday
    if days_until_saturday == 0:
        days_until_saturday = 7
    
    # Add the weeks ahead
    days_to_add = days_until_saturday + (weeks_ahead * 7)
    
    target_date = today + datetime.timedelta(days=days_to_add)
    return target_date

def get_parasha_data(weeks_ahead=0):
    """Fetch the parasha for the specified week and return a clean list of verses."""
    try:
        # Calculate the target Shabbat date
        target_date = get_next_shabbat_date(weeks_ahead)
        year = target_date.year
        month = target_date.month
        day = target_date.day
        
        print(f"Fetching parasha for date: {year}-{month:02d}-{day:02d} (weeks_ahead: {weeks_ahead})")
        
        # We specify the desired translation version.
        text_version = "vtitle=The_Contemporary_Torah,_JPS,_2006"
        
        # Use year, month, and day parameters as per Sefaria docs
        initial_cal_url = f"https://www.sefaria.org/api/calendars?year={year}&month={month}&day={day}"
        cal = requests.get(initial_cal_url).json()
        parasha_item = next(i for i in cal["calendar_items"]
                            if i["title"]["en"] == "Parashat Hashavua")

        # Extract the Gregorian date of the Parasha
        parasha_gregorian_date = cal.get("date")

        # Get Hebrew date from the same API call
        correct_hebrew_date = cal.get("hebrewDateStr", "")

        ref = parasha_item["ref"]
        print(f"Fetching parasha for date: {year}-{month:02d}-{day:02d} (weeks_ahead: {weeks_ahead})")
        
        # Parse the reference to get the correct chapter range
        # Format: "Numbers 16:1-18:32"
        ref_match = re.match(r'(\w+)\s+(\d+):(\d+)-(\d+):(\d+)', ref)
        if ref_match:
            book, start_chapter, start_verse, end_chapter, end_verse = ref_match.groups()
            start_chapter, start_verse, end_chapter, end_verse = map(int, [start_chapter, start_verse, end_chapter, end_verse])
        else:
            start_chapter = end_chapter = 1
        
        # We fetch the raw text now and will clean it ourselves.
        text_url = f"https://www.sefaria.org/api/texts/{ref}?{text_version}&context=0"
        text_data = requests.get(text_url).json()

        all_verses = []
        # Use the correctly parsed chapter range instead of API sections
        expected_chapters = list(range(start_chapter, end_chapter + 1))
        
        # Process all available text data instead of relying on sections array
        text_chapters = text_data.get('text', [])
        hebrew_chapters = text_data.get('he', [])
        section_names = text_data.get('sectionNames', [])
        
        # Process each chapter of text data
        for chap_idx in range(len(text_chapters)):
            if chap_idx >= len(expected_chapters):
                continue
                
            correct_chapter = expected_chapters[chap_idx]
            
            if chap_idx >= len(hebrew_chapters):
                continue

            en_verses_for_chap = text_chapters[chap_idx]
            he_verses_for_chap = hebrew_chapters[chap_idx]

            try:
                if chap_idx < len(section_names):
                    start_vs = int(section_names[chap_idx].split(':')[1])
                else:
                    start_vs = 1
            except (IndexError, ValueError):
                start_vs = 1

            for verse_idx, (en, he) in enumerate(zip(en_verses_for_chap, he_verses_for_chap)):
                verse_num = start_vs + verse_idx
                all_verses.append({
                    "chapter": correct_chapter,
                    "verse": verse_num,
                    "en": clean_text(en),
                    "he": clean_text(he)
                })

        if not all_verses:
             raise ValueError("No verses were parsed. Check API response.")

        return {
            "title_en": parasha_item["displayValue"]["en"],
            "hebrew_date": correct_hebrew_date,
            "gregorian_date": target_date.strftime("%A, %B %d, %Y"),
            "parasha_ref": parasha_item["ref"],
            "book": text_data["book"],
            "verses": all_verses
        }

    except Exception as e:
        print(f"An error occurred in get_parasha_data: {e}")
        return None

def create_presentation(data, output, verse_ranges=None):
    """Generates a PPTX file and saves it to the given output (path or stream)."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]
    verses_per_slide = 5

    # If verse_ranges is provided, group verses by range
    if verse_ranges:
        all_ranges = data.get('ranges')
        if not all_ranges:
            flat_verses = data['verses']
            ranges = [r.strip() for r in verse_ranges.split(',')]
            idx = 0
            all_ranges = []
            for range_str in ranges:
                parts = range_str.split(':')
                if len(parts) == 2:
                    chapter = int(parts[0])
                    if '-' in parts[1]:
                        start_verse, end_verse = map(int, parts[1].split('-'))
                    else:
                        start_verse = end_verse = int(parts[1])
                    count = end_verse - start_verse + 1
                else:
                    count = verses_per_slide
                verses = flat_verses[idx:idx+count]
                all_ranges.append({'range': range_str, 'book': data.get('book', ''), 'verses': verses})
                idx += count
    else:
        all_ranges = [{'range': None, 'book': data.get('book', ''), 'verses': data['verses']}]

    for group in all_ranges:
        verses = group['verses']
        book_name = group.get('book', data.get('book', ''))
        
        # Group verses by chapter to ensure each chapter starts on a new slide
        current_chapter_verses = []
        current_chapter = None
        
        for verse in verses:
            # If we're starting a new chapter, create a slide for the previous chapter
            if current_chapter is not None and verse['chapter'] != current_chapter:
                # Create slide(s) for the previous chapter
                for i in range(0, len(current_chapter_verses), verses_per_slide):
                    verse_chunk = current_chapter_verses[i:i+verses_per_slide]
                    slide = prs.slides.add_slide(blank_layout)
                    
                    # Title for this chunk
                    if len(verse_chunk) == 1:
                        v = verse_chunk[0]
                        title_text = f"{book_name} {v['chapter']}:{v['verse']}"
                    else:
                        start = verse_chunk[0]
                        end = verse_chunk[-1]
                        title_text = f"{book_name} {start['chapter']}:{start['verse']}-{end['verse']}"
                    
                    # Add title and content to slide
                    add_content_to_slide(slide, title_text, verse_chunk, book_name)
                
                # Start new chapter
                current_chapter_verses = [verse]
                current_chapter = verse['chapter']
            else:
                # Same chapter, add to current group
                if current_chapter is None:
                    current_chapter = verse['chapter']
                current_chapter_verses.append(verse)
        
        # Don't forget the last chapter
        if current_chapter_verses:
            for i in range(0, len(current_chapter_verses), verses_per_slide):
                verse_chunk = current_chapter_verses[i:i+verses_per_slide]
                slide = prs.slides.add_slide(blank_layout)
                
                # Title for this chunk
                if len(verse_chunk) == 1:
                    v = verse_chunk[0]
                    title_text = f"{book_name} {v['chapter']}:{v['verse']}"
                else:
                    start = verse_chunk[0]
                    end = verse_chunk[-1]
                    title_text = f"{book_name} {start['chapter']}:{start['verse']}-{end['verse']}"
                
                # Add title and content to slide
                add_content_to_slide(slide, title_text, verse_chunk, book_name)

    prs.save(output)

def add_content_to_slide(slide, title_text, verse_chunk, book_name):
    """Add title and content to a slide"""
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), width=Inches(12.333), height=Inches(0.75))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = title_text
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Sylfaen'
    p.font.size = Pt(28)
    p.font.bold = True

    # Add content
    en_parts = []
    he_parts = []
    for j, verse in enumerate(verse_chunk):
        if j > 0:
            prev_verse = verse_chunk[j-1]
            if (prev_verse['chapter'] == verse['chapter'] and verse['verse'] != prev_verse['verse'] + 1):
                en_parts.append("...")
                he_parts.append("...")
            elif prev_verse['chapter'] != verse['chapter']:
                en_parts.append("...")
                he_parts.append("...")
        en_parts.append(verse['en'])
        he_parts.append(verse['he'])
    en_text = " ".join(en_parts)
    he_text = " ".join(he_parts)

    # English text
    en_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), width=Inches(6.0), height=Inches(5.8))
    tf_en = en_box.text_frame
    tf_en.word_wrap = True
    p_en = tf_en.add_paragraph()
    p_en.text = en_text
    p_en.font.name = 'Sylfaen'
    p_en.font.size = Pt(20)

    # Hebrew text
    he_box = slide.shapes.add_textbox(Inches(6.833), Inches(1.2), width=Inches(6.0), height=Inches(5.8))
    tf_he = he_box.text_frame
    tf_he.word_wrap = True
    p_he = tf_he.paragraphs[0]
    p_he.alignment = PP_ALIGN.RIGHT
    run = p_he.add_run()
    run.text = he_text
    font = run.font
    font.name = 'Times New Roman'
    font.size = Pt(30)
    font.bold = True

if __name__ == "__main__":
    print("Fetching weekly Parasha from Sefaria.org...")
    parasha_data = get_parasha_data()
    if parasha_data:
        print(f"Found: {parasha_data['title_en']}")
        print(f"Total verses found: {len(parasha_data['verses'])}")
        file_name = f"{parasha_data['title_en']}.pptx"
        print(f"Generating PowerPoint presentation: {file_name}")
        create_presentation(parasha_data, output=file_name)
        print("Presentation saved successfully.")
